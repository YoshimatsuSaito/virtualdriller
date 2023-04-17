from pptx import Presentation


def extract_text_from_ppt(uploaded_ppt):
    """Extract text from a PowerPoint file."""
    presentation = Presentation(uploaded_ppt)
    extracted_text = ""

    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    texts.append(paragraph.text)
        extracted_text += "\n".join(texts)

    return extracted_text